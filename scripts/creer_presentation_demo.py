"""
Génère la présentation PowerPoint de démo Qual.IA (14 slides).
Thème : fond sombre #0F172A, bleu #3B82F6, vert #10B981, texte blanc.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─── Palette ─────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0F, 0x17, 0x2A)   # navy foncé
BLUE     = RGBColor(0x3B, 0x82, 0xF6)   # bleu électrique
GREEN    = RGBColor(0x10, 0xB9, 0x81)   # vert
GRAY     = RGBColor(0x94, 0xA3, 0xB8)   # gris clair
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG  = RGBColor(0x1E, 0x29, 0x3B)   # fond carte légèrement plus clair
ORANGE   = RGBColor(0xF5, 0x9E, 0x0B)

W = Inches(10)
H = Inches(5.625)

# ─── Helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def add_slide(prs):
    blank = prs.slide_layouts[6]   # blank layout
    slide = prs.slides.add_slide(blank)
    # fond uni
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def txt(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.auto_size = None
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name = "Calibri"
    return txBox


def rect(slide, x, y, w, h, fill_color, alpha=None):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def bullet_block(slide, items, x, y, w, h, size=15, color=WHITE,
                 indent_color=None, bold_first=False):
    """Ajoute un bloc de bullets. items = liste de (texte, niveau_indent)"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.auto_size = None
    tf.word_wrap = True
    first = True
    for i, item in enumerate(items):
        if isinstance(item, str):
            text, indent = item, 0
        else:
            text, indent = item

        p = tf.paragraphs[0] if (i == 0) else tf.add_paragraph()
        p.level = indent
        p.space_before = Pt(3)
        if indent == 0:
            p.space_before = Pt(6)

        run = p.add_run()
        run.text = ("• " if indent == 0 else "  — ") + text
        run.font.size = Pt(size if indent == 0 else size - 1)
        run.font.color.rgb = color if indent == 0 else GRAY
        run.font.bold = (bold_first and first)
        run.font.name = "Calibri"
        first = False
    return txBox


def add_title(slide, title_text, subtitle_text=None, caption=None):
    """Titre principal centré avec accent bleu sur le mot clé."""
    # Bande accent en haut
    rect(slide, 0, 0, 10, 0.06, BLUE)

    txt(slide, title_text, 0.5, 0.18, 9, 0.9,
        size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    if subtitle_text:
        txt(slide, subtitle_text, 0.5, 1.05, 9, 0.5,
            size=18, color=GRAY, align=PP_ALIGN.LEFT)

    if caption:
        txt(slide, caption, 0.5, 5.1, 9, 0.4,
            size=11, color=GRAY, align=PP_ALIGN.LEFT, italic=True)

    # Ligne séparatrice sous le titre
    sep = slide.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(3), Inches(0.03))
    sep.fill.solid()
    sep.fill.fore_color.rgb = BLUE
    sep.line.fill.background()


def badge(slide, text, x, y, color=BLUE):
    """Petit badge coloré."""
    r = rect(slide, x, y, len(text) * 0.12 + 0.3, 0.32, color)
    txt(slide, text, x + 0.08, y + 0.03, len(text) * 0.12 + 0.3, 0.28,
        size=11, bold=True, color=WHITE)


def card(slide, x, y, w, h, title=None, title_color=BLUE):
    """Carte avec fond légèrement différent."""
    r = rect(slide, x, y, w, h, CARD_BG)
    if title:
        txt(slide, title, x + 0.18, y + 0.12, w - 0.3, 0.35,
            size=13, bold=True, color=title_color)
    return r


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

def slide_01_titre(prs):
    """Slide 1 — Titre hero"""
    s = add_slide(prs)

    # Grand rectangle accent diagonal
    big = rect(s, 5.5, 0, 4.5, 5.625, BLUE)
    big.fill.fore_color.rgb = RGBColor(0x1E, 0x40, 0xAF)

    # Logo / icone horlogerie
    txt(s, "⌚", 6.2, 0.8, 3, 1.5, size=72, align=PP_ALIGN.CENTER, color=WHITE)

    # Titre
    txt(s, "Qual", 0.7, 1.1, 5, 1.1, size=64, bold=True, color=WHITE)
    txt(s, ".IA", 0.7, 1.1, 5, 1.1, size=64, bold=True, color=BLUE)

    # Hack : superposer ".IA" en bleu sur "Qual" en blanc
    # Réécrivons en 2 textbox côte à côte
    # Supprimer les 2 précédents n'est pas possible, on les réécrit correctement :
    for sh in list(s.shapes):
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if run.text in ("Qual", ".IA"):
                        sh.element.getparent().remove(sh.element)
                        break

    txt(s, "Qual.IA", 0.7, 1.1, 5, 1.1, size=60, bold=True, color=WHITE)

    txt(s, "Génération automatique de documents\nqualité industriels par l'IA",
        0.7, 2.35, 5, 1.2, size=19, color=GRAY)

    txt(s, "MVP v1.1  ·  Horlogerie  ·  Céramique  ·  Multi-secteur",
        0.7, 4.6, 5, 0.5, size=12, color=GRAY, italic=True)

    # Tags
    badge(s, "Claude API", 0.7, 3.8, BLUE)
    badge(s, "Streamlit", 2.2, 3.8, RGBColor(0x0D, 0x94, 0x88))
    badge(s, "ISO 9001", 3.7, 3.8, GREEN)

    rect(s, 0, 0, 10, 0.06, BLUE)


def slide_02_probleme(prs):
    s = add_slide(prs)
    add_title(s, "Le problème aujourd'hui", "Rédiger des documents qualité est lent, risqué et non traçable")
    rect(s, 0, 0, 10, 0.06, ORANGE)

    problems = [
        "Rédiger une AMDEC Produit + Process + Gamme = 2 à 3 jours de travail",
        "Risque d'erreur humaine élevé (copier-coller entre fichiers Excel)",
        "Chaque nouveau produit recommence from scratch — aucune capitalisation",
        "Pas de workflow de validation : qui a approuvé ? quand ?",
        "La connaissance qualité est bloquée dans des fichiers individuels",
    ]
    for i, p in enumerate(problems):
        y = 1.45 + i * 0.75
        # Cercle numéro
        circle = s.shapes.add_shape(9, Inches(0.5), Inches(y), Inches(0.42), Inches(0.42))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ORANGE
        circle.line.fill.background()
        txt(s, str(i+1), 0.5, y - 0.01, 0.42, 0.42,
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, p, 1.1, y + 0.01, 8.5, 0.5, size=14, color=WHITE)


def slide_03_solution(prs):
    s = add_slide(prs)
    add_title(s, "Qual.IA — La solution", "L'IA adapte des documents validés. Elle n'invente pas.")
    rect(s, 0, 0, 10, 0.06, GREEN)

    items = [
        ("L'IA adapte les documents d'une référence validée existante", GREEN),
        ("AMDEC Produit + AMDEC Process + Gamme générés en 30 secondes", BLUE),
        ("Aucune invention : l'IA adapte uniquement ce qui change (cotes, traitements)", WHITE),
        ("Base de connaissances qui s'enrichit à chaque produit validé", GREEN),
        ("Workflow intégré : BT → Qualité → Direction", BLUE),
    ]
    for i, (txt_val, col) in enumerate(items):
        y = 1.45 + i * 0.77
        # Check ou flèche
        icon = "✓" if col == GREEN else "→"
        ic_col = col
        txt(s, icon, 0.5, y, 0.5, 0.5, size=18, bold=True, color=ic_col)
        txt(s, txt_val, 1.1, y + 0.03, 8.5, 0.5, size=14, color=WHITE)


def slide_04_archi(prs):
    s = add_slide(prs)
    add_title(s, "Architecture technique", "Stack léger, déployable localement ou en cloud")

    # Colonne gauche — Frontend
    card(s, 0.4, 1.25, 4.4, 3.8, "🖥  Frontend — Streamlit", BLUE)
    fe_items = [
        "Interface web locale (navigateur)",
        "4 pages de navigation",
        "  Nouveau produit / Import Excel",
        "  Workflow / Gestion de la base",
        "Tableaux éditables cellule par cellule",
        "Upload PDF (plan technique)",
    ]
    for i, item in enumerate(fe_items):
        indent = item.startswith("  ")
        txt(s, item.strip(), 0.65, 1.75 + i * 0.47, 4.0, 0.45,
            size=12 if not indent else 11,
            color=WHITE if not indent else GRAY)

    # Colonne droite — Backend
    card(s, 5.2, 1.25, 4.4, 3.8, "⚙  Backend — Python", GREEN)
    be_items = [
        "Claude API (claude-sonnet-4-6)",
        "  3 documents en 1 appel groupé",
        "Claude Vision — lecture plans PDF",
        "Stockage JSON (pas de base SQL)",
        "Sauvegarde automatique quotidienne",
        "Multi-catégories isolées",
    ]
    for i, item in enumerate(be_items):
        indent = item.startswith("  ")
        txt(s, item.strip(), 5.45, 1.75 + i * 0.47, 4.0, 0.45,
            size=12 if not indent else 11,
            color=WHITE if not indent else GRAY)

    # Flèche centrale
    txt(s, "⇄", 4.6, 2.8, 0.7, 0.6, size=28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


def slide_05_flux(prs):
    s = add_slide(prs)
    add_title(s, "Comment ça marche", "6 étapes du brief client au document validé")
    rect(s, 0, 0, 10, 0.06, BLUE)

    steps = [
        ("①", "Brief client",      "Diamètre, matière, traitements, tolérances"),
        ("②", "Import plan PDF",   "Claude Vision pré-remplit le formulaire"),
        ("③", "Similarité",        "Trouve la meilleure référence source (score %)"),
        ("④", "Génération IA",     "AMDEC Produit + Process + Gamme en 30 sec"),
        ("⑤", "Révision",          "Édition cellule par cellule dans l'interface"),
        ("⑥", "Sauvegarde",        "Devient nouvelle référence dans la base"),
    ]
    cols = [(0.35, BLUE), (3.6, GREEN), (6.85, BLUE)]
    rows = [(0, 1), (2, 3), (4, 5)]

    for (col_x, col_color), (idx1, idx2) in zip(cols, rows):
        for row_offset, step_idx in enumerate(range(idx1, idx2 + 1)):
            if step_idx >= len(steps):
                break
            num, title, desc = steps[step_idx]
            y = 1.35 + row_offset * 1.85
            card(s, col_x, y, 3.0, 1.65)
            txt(s, num, col_x + 0.15, y + 0.12, 0.55, 0.55,
                size=22, bold=True, color=col_color)
            txt(s, title, col_x + 0.72, y + 0.12, 2.15, 0.45,
                size=13, bold=True, color=WHITE)
            txt(s, desc, col_x + 0.15, y + 0.65, 2.7, 0.75,
                size=11, color=GRAY)


def slide_06_similarite(prs):
    s = add_slide(prs)
    add_title(s, "Moteur de similarité", "Trouve automatiquement la meilleure référence source")

    # Score pondéré
    weights = [
        ("Diamètre", 40, BLUE),
        ("Matière", 25, GREEN),
        ("Traitements", 20, ORANGE),
        ("Épaisseur", 10, GRAY),
        ("Tolérance", 5, GRAY),
    ]
    for i, (label, pct, col) in enumerate(weights):
        x = 0.5 + i * 1.88
        # Barre
        full_h = 1.6
        bar_h = full_h * pct / 40
        rect(s, x, 1.4 + (full_h - bar_h), 1.3, bar_h, col)
        txt(s, f"{pct}%", x, 1.2, 1.3, 0.4, size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(s, label, x, 3.06, 1.3, 0.4, size=12, color=GRAY, align=PP_ALIGN.CENTER)

    # 3 modes
    modes = [
        ("> 85%", "Mode automatique", GREEN, "✅"),
        ("60–85%", "Avertissement", ORANGE, "⚠"),
        ("< 60%", "Choix manuel", GRAY, "🔧"),
    ]
    for i, (score, label, col, icon) in enumerate(modes):
        x = 0.5 + i * 3.15
        card(s, x, 3.7, 2.9, 1.55, title_color=col)
        txt(s, icon + " " + score, x + 0.18, 3.82, 2.5, 0.45, size=16, bold=True, color=col)
        txt(s, label, x + 0.18, 4.28, 2.5, 0.75, size=13, color=WHITE)


def slide_07_variantes(prs):
    s = add_slide(prs)
    add_title(s, "Génération de variantes", "Même produit, traitements différents → N articles en 1 fois")

    # Cas d'usage
    txt(s, "Cas d'usage :", 0.5, 1.35, 4, 0.4, size=13, bold=True, color=GRAY)

    articles = [
        ("MA696.361", "Glace Ø31mm", "Métallisation CN", BLUE),
        ("MB696.361", "Glace Ø31mm", "Métallisation AU", GREEN),
        ("MC696.361", "Glace Ø31mm", "Métallisation CC", ORANGE),
    ]
    for i, (art, base, trait, col) in enumerate(articles):
        y = 1.7 + i * 0.9
        card(s, 0.5, y, 4.3, 0.78)
        txt(s, art, 0.75, y + 0.08, 1.6, 0.45, size=13, bold=True, color=col)
        txt(s, base, 2.45, y + 0.08, 1.4, 0.35, size=11, color=WHITE)
        txt(s, "+ " + trait, 2.45, y + 0.38, 2.0, 0.35, size=11, color=GRAY)

    # Explication droite
    card(s, 5.2, 1.35, 4.4, 3.8, "💡 Comment faire", BLUE)
    exp_items = [
        "Dans le brief client, cocher",
        "  « Générer plusieurs variantes »",
        "Ajouter 1 ligne par variante :",
        "  N° article + traitements",
        "  Désignation client (optionnel)",
        "3 appels API pour N variantes",
        "  (pas N×3 — optimisation ÷N)",
        "Export Excel de tous les dossiers",
    ]
    for i, item in enumerate(exp_items):
        indent = item.startswith("  ")
        txt(s, item.strip(), 5.45, 1.85 + i * 0.42, 3.9, 0.4,
            size=12 if not indent else 11,
            color=WHITE if not indent else GRAY)


def slide_08_vision(prs):
    s = add_slide(prs)
    add_title(s, "Import de plan technique", "Claude Vision lit le plan et pré-remplit le formulaire")

    # Côté gauche : représentation plan
    card(s, 0.4, 1.3, 4.5, 3.8)
    txt(s, "📐", 0.9, 1.55, 1.5, 1.0, size=48, align=PP_ALIGN.CENTER, color=BLUE)
    txt(s, "PDF / PNG / JPG", 0.5, 2.55, 4.2, 0.4, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, "Plan technique industriel", 0.5, 2.95, 4.2, 0.4, size=12, color=GRAY, align=PP_ALIGN.CENTER)

    extracted = [
        "Diamètre + épaisseur (± tolérances)",
        "Type de produit et matière",
        "Traitements de surface (CN, AR, PVD…)",
        "Désignation du cartouche",
    ]
    for i, e in enumerate(extracted):
        txt(s, "↳ " + e, 0.65, 3.5 + i * 0.38, 4.1, 0.35, size=11, color=GRAY)

    # Flèche
    txt(s, "→", 4.6, 2.9, 0.7, 0.5, size=28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    txt(s, "1 clic", 4.55, 3.45, 0.8, 0.4, size=10, color=GRAY, align=PP_ALIGN.CENTER)

    # Côté droit : formulaire pré-rempli
    card(s, 5.2, 1.3, 4.4, 3.8, "Formulaire pré-rempli", GREEN)
    fields = [
        ("Diamètre", "31.00 mm"),
        ("Tolérance Ø", "± 0.05 mm"),
        ("Épaisseur", "2.50 mm"),
        ("Matière", "Saphir"),
        ("Type", "Glace ronde"),
        ("Traitements", "CN + AR double face"),
    ]
    for i, (label, val) in enumerate(fields):
        y = 1.82 + i * 0.51
        txt(s, label + " :", 5.45, y, 1.7, 0.4, size=11, color=GRAY)
        txt(s, val, 7.2, y, 2.2, 0.4, size=11, bold=True, color=WHITE)


def slide_09_workflow(prs):
    s = add_slide(prs)
    add_title(s, "Workflow de validation qualité", "Traçabilité complète de la création à la libération")
    rect(s, 0, 0, 10, 0.06, GREEN)

    # Diagramme états
    states = [
        ("📝", "Brouillon",   GRAY,  0.45),
        ("👁",  "En revue",    BLUE,  2.75),
        ("✅", "Approuvé",    GREEN, 5.05),
        ("🚀", "Libéré",      GREEN, 7.35),
    ]
    for icon, label, col, x in states:
        circle = s.shapes.add_shape(9, Inches(x), Inches(1.32), Inches(1.3), Inches(1.3))
        circle.fill.solid()
        circle.fill.fore_color.rgb = col
        circle.line.fill.background()
        txt(s, icon, x, 1.32, 1.3, 0.8, size=28, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, label, x - 0.1, 2.62, 1.5, 0.4, size=11, bold=True, color=col, align=PP_ALIGN.CENTER)

    # Flèches entre états
    for fx in [1.8, 4.1, 6.4]:
        txt(s, "→", fx, 1.65, 0.85, 0.6, size=22, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

    # Boucle corrections
    txt(s, "🔄 Corrections demandées", 2.4, 3.2, 3.5, 0.5, size=12, color=ORANGE, align=PP_ALIGN.CENTER)
    txt(s, "↕", 3.08, 2.9, 0.5, 0.4, size=18, color=ORANGE, align=PP_ALIGN.CENTER)

    # Gates tableau
    gates = [
        ("IPR ≤ 40",   "Gate BT uniquement",                      GRAY),
        ("IPR 41–100", "BT + Approbation Qualité",                 BLUE),
        ("IPR > 100",  "BT + Qualité + Validation Direction",      ORANGE),
    ]
    txt(s, "Gates selon IPR (G × O × D) :", 0.5, 3.75, 5, 0.4, size=12, bold=True, color=GRAY)
    for i, (ipr, gates_txt, col) in enumerate(gates):
        y = 4.15 + i * 0.38
        txt(s, ipr, 0.6, y, 1.4, 0.35, size=11, bold=True, color=col)
        txt(s, "→  " + gates_txt, 2.1, y, 7.4, 0.35, size=11, color=WHITE)


def slide_10_dit(prs):
    s = add_slide(prs)
    add_title(s, "DIT — Instructions de Travail", "Un document par opération, généré par l'IA")

    # Schéma liaison Gamme → DIT
    txt(s, "GAMME", 0.5, 1.35, 2.0, 0.55, size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    rect(s, 0.5, 1.35, 2.0, 0.55, CARD_BG)
    txt(s, "GAMME", 0.5, 1.35, 2.0, 0.55, size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    ops = [("Op 10", "Réception"), ("Op 30", "Nettoyage"), ("Op 50", "Métallisation")]
    for i, (no, name) in enumerate(ops):
        y = 2.05 + i * 0.88
        rect(s, 0.5, y, 1.95, 0.65, CARD_BG)
        txt(s, no, 0.65, y + 0.05, 0.7, 0.35, size=11, bold=True, color=BLUE)
        txt(s, name, 1.4, y + 0.05, 0.95, 0.55, size=10, color=WHITE)
        # Flèche
        txt(s, "→", 2.5, y + 0.12, 0.5, 0.35, size=16, color=GRAY)
        # DIT
        dit_code = f"DIT-{(i+1)*10:03d}"
        rect(s, 3.05, y, 1.6, 0.65, CARD_BG)
        txt(s, dit_code, 3.1, y + 0.05, 1.5, 0.55, size=11, bold=True, color=GREEN)

    # Contenu d'un DIT
    card(s, 5.0, 1.3, 4.6, 4.0, "📄 Contenu d'un DIT", GREEN)
    dit_content = [
        "Code et titre de l'instruction",
        "Révision + Rédigé par",
        "EPI requis (gants, lunettes…)",
        "Outillage nécessaire",
        "Étapes numérotées avec",
        "  points de vigilance",
        "Critères d'acceptation",
        "Notes de sécurité",
    ]
    for i, item in enumerate(dit_content):
        indent = item.startswith("  ")
        txt(s, ("  " if indent else "• ") + item.strip(), 5.25, 1.82 + i * 0.44, 4.15, 0.42,
            size=12 if not indent else 11,
            color=WHITE if not indent else GRAY)

    # Badge IA
    badge(s, "🤖 Généré par IA en 1 clic", 0.5, 5.05, GREEN)


def slide_11_base(prs):
    s = add_slide(prs)
    add_title(s, "Gestion de la base de références", "Toutes vos références qualité en un seul endroit")

    # Métriques
    metrics = [
        ("3", "Catégories", BLUE),
        ("14", "Références", GREEN),
        ("5", "Onglets/réf", ORANGE),
    ]
    for i, (val, label, col) in enumerate(metrics):
        x = 0.5 + i * 3.15
        card(s, x, 1.3, 2.9, 1.3)
        txt(s, val, x + 0.2, 1.35, 1.2, 0.9, size=48, bold=True, color=col)
        txt(s, label, x + 1.45, 1.6, 1.3, 0.6, size=14, color=GRAY)

    # Onglets
    tabs = [
        ("📋 Métadonnées", "Statut, approbateur, dimensions, workflow"),
        ("AMDEC Produit",  "Modes de défaillance éditable + IPR calculé"),
        ("AMDEC Process",  "Défaillances process, paramètres clés"),
        ("Gamme",          "Opérations + temps total calculé automatiquement"),
        ("📄 DIT",         "Instructions de travail par opération"),
        ("⚙ Actions",      "Modifier statut, supprimer avec sauvegarde"),
    ]
    for i, (tab, desc) in enumerate(tabs):
        col_i = i % 2
        row_i = i // 2
        x = 0.5 + col_i * 4.85
        y = 2.9 + row_i * 0.78
        rect(s, x, y, 0.08, 0.5, BLUE)
        txt(s, tab, x + 0.18, y + 0.02, 2.5, 0.35, size=12, bold=True, color=WHITE)
        txt(s, desc, x + 0.18, y + 0.32, 4.3, 0.35, size=10, color=GRAY)


def slide_12_multicategories(prs):
    s = add_slide(prs)
    add_title(s, "Architecture multi-catégories", "Chaque secteur métier est un espace isolé")

    cats = [
        ("⌚", "Couches\nMinces",  "Glace saphir + dépôts\nCN, AR, PVD, DLC…", BLUE),
        ("🏺", "Glace\nNue",      "Glace saphir/minéral\nsans traitement", GREEN),
        ("⚙", "Céramique",        "Pièces céramique\nZircone, Alumine…", ORANGE),
        ("✈", "Aéronautique\n(à venir)", "Pièces métal/composite\nAMS, NADCAP…", GRAY),
    ]
    for i, (icon, name, desc, col) in enumerate(cats):
        x = 0.4 + i * 2.38
        card(s, x, 1.35, 2.2, 2.5)
        txt(s, icon, x + 0.1, 1.45, 1.9, 0.8, size=36, align=PP_ALIGN.CENTER, color=col)
        txt(s, name, x + 0.1, 2.2, 1.9, 0.65, size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(s, desc, x + 0.1, 2.85, 1.9, 0.75, size=10, color=GRAY, align=PP_ALIGN.CENTER)

    # config.json
    card(s, 0.4, 4.05, 9.2, 1.28, "config.json de la catégorie :", GRAY)
    config_items = "Critères similarité & poids  ·  Vocabulaire (types, matières, traitements)  ·  Prompt expert Claude"
    txt(s, config_items, 0.65, 4.48, 8.9, 0.5, size=12, color=WHITE)
    txt(s, "Ajouter une catégorie = créer un dossier + 1 fichier config.json", 0.65, 4.82, 8.9, 0.4,
        size=11, color=GRAY, italic=True)


def slide_13_roadmap(prs):
    s = add_slide(prs)
    add_title(s, "Roadmap", "Ce qui est fait et ce qui arrive")

    # Colonne Fait
    card(s, 0.4, 1.3, 4.3, 4.0, "✅  Déjà disponible", GREEN)
    done = [
        "Génération AMDEC + Gamme par IA",
        "Variantes multi-articles (3 appels pour N)",
        "Import plan technique Vision",
        "Workflow de validation avec gates",
        "DIT par opération (génération IA)",
        "Import Excel client",
        "Sauvegarde automatique quotidienne",
    ]
    for i, item in enumerate(done):
        txt(s, "✓  " + item, 0.65, 1.82 + i * 0.47, 3.9, 0.42, size=11.5, color=WHITE)

    # Colonne Prévu
    card(s, 5.3, 1.3, 4.3, 4.0, "🚧  Prochaines étapes", ORANGE)
    todo = [
        "Export DIT en PDF imprimable",
        "Notifications email (workflow)",
        "Multi-tenant (plusieurs entreprises)",
        "Connecteur ERP SAP / Infor",
        "Produits composites automatiques",
        "Tableau de bord analytique",
        "Application mobile (consultation)",
    ]
    for i, item in enumerate(todo):
        txt(s, "→  " + item, 5.55, 1.82 + i * 0.47, 3.9, 0.42, size=11.5, color=WHITE)


def slide_14_contact(prs):
    s = add_slide(prs)

    # Fond avec dégradé simulé par rectangle
    rect(s, 0, 0, 10, 5.625, CARD_BG)
    rect(s, 0, 0, 10, 0.06, GREEN)

    txt(s, "⌚  Qual.IA", 0.7, 0.55, 9, 0.85,
        size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(s, "Automatisez votre qualité. Gardez l'expertise.",
        0.7, 1.45, 9, 0.65, size=24, color=GREEN, align=PP_ALIGN.CENTER, italic=True)

    infos = [
        "🤖  Développé avec Claude API (Anthropic)",
        "💼  Logiciel SaaS — licence mensuelle par entreprise",
        "🌍  Compatible tous secteurs industriels",
        "📅  Démo disponible sur rendez-vous",
    ]
    for i, info in enumerate(infos):
        txt(s, info, 1.5, 2.45 + i * 0.55, 7, 0.48, size=15, color=WHITE, align=PP_ALIGN.CENTER)

    badges = [
        ("ISO 9001", GREEN),
        ("AMDEC", BLUE),
        ("Gamme", BLUE),
        ("DIT", GREEN),
        ("Workflow", ORANGE),
    ]
    for i, (label, col) in enumerate(badges):
        x = 1.4 + i * 1.55
        badge(s, label, x, 4.8, col)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    slide_01_titre(prs)
    slide_02_probleme(prs)
    slide_03_solution(prs)
    slide_04_archi(prs)
    slide_05_flux(prs)
    slide_06_similarite(prs)
    slide_07_variantes(prs)
    slide_08_vision(prs)
    slide_09_workflow(prs)
    slide_10_dit(prs)
    slide_11_base(prs)
    slide_12_multicategories(prs)
    slide_13_roadmap(prs)
    slide_14_contact(prs)

    out = r"C:\Users\yadel\Desktop\Projet dev SaaS\projet-qualite-ia\DEMO_QUALIA.pptx"
    prs.save(out)
    print(f"✅ Présentation créée : {out}")
    print(f"   {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
