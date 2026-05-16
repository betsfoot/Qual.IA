"""
Génère une présentation PowerPoint de pitch pour Qual.IA.
Lancer : python scripts/generer_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pathlib import Path
import sys

# ─── Palette couleurs ────────────────────────────────────────────────────────
BLEU_FONCE  = RGBColor(0x0D, 0x1B, 0x3E)   # fond slides principales
BLEU_VIF    = RGBColor(0x1A, 0x6B, 0xFF)   # accents, titres
BLANC       = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_CLAIR  = RGBColor(0xF2, 0xF4, 0xF8)
VERT        = RGBColor(0x00, 0xC2, 0x7A)
ORANGE      = RGBColor(0xFF, 0x6B, 0x2B)
GRIS_TEXTE  = RGBColor(0x44, 0x44, 0x55)

OUTPUT = Path(__file__).parent.parent / "exports" / "QualIA_Pitch.pptx"
OUTPUT.parent.mkdir(exist_ok=True)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ─── Helpers ─────────────────────────────────────────────────────────────────

def blank_slide():
    layout = prs.slide_layouts[6]  # totalement vide
    return prs.slides.add_slide(layout)

def rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def txt(slide, text, x, y, w, h, size=18, bold=False, color=BLANC,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def badge(slide, label, x, y, w=Inches(2.2), h=Inches(0.45), bg=BLEU_VIF, fg=BLANC, size=13):
    rect(slide, x, y, w, h, bg)
    txt(slide, label, x, y + Inches(0.05), w, h, size=size, bold=True, color=fg, align=PP_ALIGN.CENTER)

def kpi(slide, valeur, label, x, y, color=BLEU_VIF):
    rect(slide, x, y, Inches(2.8), Inches(1.5), GRIS_CLAIR)
    txt(slide, valeur, x, y + Inches(0.1), Inches(2.8), Inches(0.8),
        size=32, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(slide, label, x, y + Inches(0.85), Inches(2.8), Inches(0.55),
        size=12, color=GRIS_TEXTE, align=PP_ALIGN.CENTER)

def fleche(slide, x, y):
    shape = slide.shapes.add_shape(13, x, y, Inches(0.4), Inches(0.3))  # right arrow
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLEU_VIF
    shape.line.fill.background()

def fond_bleu(slide):
    rect(slide, 0, 0, W, H, BLEU_FONCE)

def barre_titre(slide, titre, sous_titre=""):
    rect(slide, 0, 0, W, Inches(1.3), BLEU_VIF)
    txt(slide, titre, Inches(0.4), Inches(0.12), W - Inches(0.8), Inches(0.7),
        size=28, bold=True, color=BLANC)
    if sous_titre:
        txt(slide, sous_titre, Inches(0.4), Inches(0.75), W - Inches(0.8), Inches(0.45),
            size=14, color=RGBColor(0xCC, 0xDD, 0xFF))

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Couverture
# ═════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fond_bleu(sl)
rect(sl, 0, H - Inches(0.15), W, Inches(0.15), BLEU_VIF)
rect(sl, 0, 0, Inches(0.15), H, BLEU_VIF)

txt(sl, "Qual.IA", Inches(1), Inches(1.5), Inches(11), Inches(1.5),
    size=60, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
txt(sl, "Générez vos dossiers qualité industriels en moins d'une heure",
    Inches(1), Inches(3.1), Inches(11), Inches(0.8),
    size=22, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

rect(sl, Inches(4.5), Inches(4.2), Inches(4.3), Inches(0.6), BLEU_VIF)
txt(sl, "AMDEC Produit  •  AMDEC Process  •  Gamme de Production",
    Inches(4.5), Inches(4.2), Inches(4.3), Inches(0.6),
    size=12, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

txt(sl, "Propulsé par Claude AI — Anthropic",
    Inches(1), Inches(5.5), Inches(11), Inches(0.4),
    size=13, italic=True, color=RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.CENTER)

txt(sl, "Confidentiel — Présentation commerciale",
    Inches(1), Inches(6.8), Inches(11), Inches(0.3),
    size=10, color=RGBColor(0x66, 0x77, 0x99), align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Le problème
# ═════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fond_bleu(sl)
barre_titre(sl, "Le problème", "Ce que vivent vos équipes qualité aujourd'hui")

# 3 blocs problème
problemes = [
    ("⏱️", "2 à 5 jours", "Pour créer un dossier qualité complet (AMDEC + Gamme) sur un nouveau produit"),
    ("🔁", "Travail répétitif", "80% du contenu est identique d'une référence à l'autre — copier-coller manuel risqué"),
    ("🚨", "Erreurs humaines", "Oubli d'un mode de défaillance, cote non mise à jour, perte de traçabilité"),
]
for i, (icone, titre, desc) in enumerate(problemes):
    x = Inches(0.5 + i * 4.2)
    y = Inches(1.8)
    rect(sl, x, y, Inches(3.9), Inches(3.8), RGBColor(0x1A, 0x2A, 0x5E))
    txt(sl, icone, x, y + Inches(0.3), Inches(3.9), Inches(0.8),
        size=36, align=PP_ALIGN.CENTER, color=BLANC)
    txt(sl, titre, x, y + Inches(1.1), Inches(3.9), Inches(0.6),
        size=20, bold=True, color=BLEU_VIF, align=PP_ALIGN.CENTER)
    txt(sl, desc, x + Inches(0.2), y + Inches(1.8), Inches(3.5), Inches(1.7),
        size=13, color=RGBColor(0xBB, 0xCC, 0xEE), align=PP_ALIGN.CENTER)

txt(sl, "Résultat : délais allongés, ressources mobilisées, risque qualité non maîtrisé",
    Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5),
    size=14, bold=True, italic=True, color=ORANGE, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — La solution
# ═════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fond_bleu(sl)
barre_titre(sl, "Notre solution", "Qual.IA — le copier-coller intelligent qui ne se trompe pas")

txt(sl, "Vous avez déjà des dossiers qualité validés.\nL'IA les adapte automatiquement pour chaque nouveau produit similaire.",
    Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.0),
    size=16, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

# Flux simplifié
etapes = [
    ("📋", "Brief\nclient"),
    ("🔍", "Similarité\nautomatique"),
    ("🤖", "Génération\nIA"),
    ("✏️", "Relecture\nBT"),
    ("📄", "Export\nExcel"),
]
for i, (icone, label) in enumerate(etapes):
    x = Inches(0.6 + i * 2.5)
    y = Inches(2.8)
    rect(sl, x, y, Inches(2.0), Inches(1.8), RGBColor(0x1A, 0x2A, 0x5E))
    txt(sl, icone, x, y + Inches(0.15), Inches(2.0), Inches(0.7),
        size=28, align=PP_ALIGN.CENTER, color=BLANC)
    txt(sl, label, x, y + Inches(0.9), Inches(2.0), Inches(0.7),
        size=12, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    if i < len(etapes) - 1:
        txt(sl, "→", Inches(2.55 + i * 2.5), y + Inches(0.6), Inches(0.5), Inches(0.6),
            size=22, bold=True, color=BLEU_VIF, align=PP_ALIGN.CENTER)

# Bénéfices
benefices = [
    ("< 1 heure", "au lieu de 2-5 jours", VERT),
    ("100%", "des modes conservés", BLEU_VIF),
    ("0 oubli", "traçabilité garantie", ORANGE),
]
for i, (val, label, col) in enumerate(benefices):
    x = Inches(1.0 + i * 4.0)
    kpi(sl, val, label, x, Inches(5.1), col)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Comment ça marche (workflow détaillé)
# ═════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fond_bleu(sl)
barre_titre(sl, "Comment ça marche", "Un workflow en 6 étapes, de la commande au dossier validé")

steps = [
    ("①", "Brief client",        "Type, matière, traitements,\ndimensions — 2 minutes"),
    ("②", "Analyse similarité",  "Le moteur compare au catalogue\net trouve la référence la plus proche"),
    ("③", "Génération IA",       "Claude API adapte AMDEC Produit,\nAMDEC Process et Gamme en 60s"),
    ("④", "Validation BT",       "Tableaux éditables — ajouter,\nmodifier, supprimer des lignes"),
    ("⑤", "Export Excel",        "3 fichiers Excel formatés,\nprêts à signer"),
    ("⑥", "Enrichissement base", "Le dossier validé enrichit\nla bibliothèque pour la suite"),
]
for i, (num, titre, desc) in enumerate(steps):
    col = i % 3
    row = i // 3
    x = Inches(0.4 + col * 4.3)
    y = Inches(1.6 + row * 2.5)
    rect(sl, x, y, Inches(4.0), Inches(2.1), RGBColor(0x12, 0x22, 0x4A))
    rect(sl, x, y, Inches(0.6), Inches(2.1), BLEU_VIF)
    txt(sl, num, x, y + Inches(0.55), Inches(0.6), Inches(0.8),
        size=18, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    txt(sl, titre, x + Inches(0.7), y + Inches(0.15), Inches(3.2), Inches(0.5),
        size=15, bold=True, color=BLANC)
    txt(sl, desc, x + Inches(0.7), y + Inches(0.7), Inches(3.1), Inches(1.2),
        size=12, color=RGBColor(0xBB, 0xCC, 0xEE))

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Multi-catégories
# ═════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fond_bleu(sl)
barre_titre(sl, "Un système multi-catégories", "Adaptable à chaque secteur industriel de votre entreprise")

txt(sl, "Chaque catégorie dispose de son propre vocabulaire, ses critères de similarité et son expert IA spécialisé",
    Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.5),
    size=14, color=RGBColor(0xBB, 0xCC, 0xEE), align=PP_ALIGN.CENTER)

cats = [
    ("🔭", "Couches minces", "Métallisation Chrome/Or\nAntireflet PVD/CVD\nSaphir, minéral, organique"),
    ("🔍", "Glace nue", "Rectification, polissage\nSaphir, quartz, PMMA\nRa ≤ 5nm poli optique"),
    ("🧱", "Céramique", "Alumine, Zircone Y-TZP\nFrittage, densité Archimède\nHV Vickers, ISO 6474"),
    ("💊", "Pharma", "Dispositifs médicaux\nISO 13485, GMP\nDossier technique CE"),
    ("⚙️", "Votre secteur", "Configurable en\nquelques heures\nVocabulaire sur mesure"),
]
for i, (icone, nom, desc) in enumerate(cats):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(2.1 + row * 2.35)
    bg = RGBColor(0x1A, 0x2A, 0x5E) if i < 3 else RGBColor(0x12, 0x22, 0x44)
    rect(sl, x, y, Inches(3.9), Inches(2.1), bg)
    if i == 5:
        rect(sl, x, y, Inches(3.9), Inches(0.08), ORANGE)
    txt(sl, f"{icone}  {nom}", x + Inches(0.15), y + Inches(0.15), Inches(3.6), Inches(0.5),
        size=15, bold=True, color=BLANC if i < 5 else ORANGE)
    txt(sl, desc, x + Inches(0.15), y + Inches(0.7), Inches(3.6), Inches(1.2),
        size=12, color=RGBColor(0xBB, 0xCC, 0xEE))

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Valeur ajoutée / ROI
# ═════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fond_bleu(sl)
barre_titre(sl, "Valeur ajoutée", "Ce que ça change concrètement pour votre équipe")

comparaison = [
    ("Avant", "Après", ""),
    ("2 à 5 jours par dossier", "< 1 heure", "Temps de création"),
    ("Copier-coller manuel Excel", "Adaptation IA automatique", "Méthode"),
    ("Risque d'oubli de modes", "100% des modes conservés", "Fiabilité"),
    ("Dossiers en silo (fichiers locaux)", "Base centralisée et versionée", "Traçabilité"),
    ("Formation longue nouveaux BT", "Prise en main < 1 journée", "Onboarding"),
    ("Export manuel, mise en forme", "Excel formaté en 1 clic", "Export"),
]

# En-tête tableau
rect(sl, Inches(0.4), Inches(1.5), Inches(4.5), Inches(0.45), GRIS_TEXTE)
rect(sl, Inches(4.9), Inches(1.5), Inches(4.0), Inches(0.45), RGBColor(0x8B, 0x00, 0x00))
rect(sl, Inches(8.9), Inches(1.5), Inches(4.0), Inches(0.45), VERT)
txt(sl, "Critère", Inches(0.4), Inches(1.5), Inches(4.5), Inches(0.45),
    size=13, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
txt(sl, "Avant Qual.IA", Inches(4.9), Inches(1.5), Inches(4.0), Inches(0.45),
    size=13, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
txt(sl, "Avec Qual.IA", Inches(8.9), Inches(1.5), Inches(4.0), Inches(0.45),
    size=13, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

for i, (avant, apres, critere) in enumerate(comparaison[1:]):
    y = Inches(1.95 + i * 0.72)
    bg = RGBColor(0x12, 0x22, 0x44) if i % 2 == 0 else RGBColor(0x0D, 0x1B, 0x3E)
    rect(sl, Inches(0.4), y, Inches(4.5), Inches(0.65), bg)
    rect(sl, Inches(4.9), y, Inches(4.0), Inches(0.65), bg)
    rect(sl, Inches(8.9), y, Inches(4.0), Inches(0.65), bg)
    txt(sl, critere, Inches(0.5), y + Inches(0.08), Inches(4.3), Inches(0.5),
        size=12, bold=True, color=RGBColor(0xBB, 0xCC, 0xEE))
    txt(sl, avant, Inches(5.0), y + Inches(0.08), Inches(3.8), Inches(0.5),
        size=12, color=RGBColor(0xFF, 0x99, 0x99))
    txt(sl, apres, Inches(9.0), y + Inches(0.08), Inches(3.8), Inches(0.5),
        size=12, color=VERT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Onboarding client
# ═════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fond_bleu(sl)
barre_titre(sl, "Mise en place", "Du premier contact au premier dossier généré")

phases = [
    ("Semaine 1", "Analyse", BLEU_VIF,
     ["Recueil de 5 à 10 dossiers\nqualité existants validés",
      "Identification des catégories\nproduit concernées",
      "Configuration du vocabulaire\nmétier"]),
    ("Semaine 2", "Intégration", ORANGE,
     ["Import des références\nen base (script automatisé)",
      "Paramétrage des critères\nde similarité",
      "Tests de génération\nsur cas réels"]),
    ("Semaine 3", "Formation", VERT,
     ["Formation équipe BT\n(< 1 journée)",
      "Génération des premiers\ndossiers en conditions réelles",
      "Ajustements et validation\nclient"]),
]

for i, (semaine, titre, col, points) in enumerate(phases):
    x = Inches(0.5 + i * 4.2)
    y = Inches(1.6)
    rect(sl, x, y, Inches(3.9), Inches(0.5), col)
    txt(sl, semaine, x, y, Inches(3.9), Inches(0.5),
        size=14, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    rect(sl, x, y + Inches(0.5), Inches(3.9), Inches(4.5), RGBColor(0x12, 0x22, 0x44))
    txt(sl, titre, x + Inches(0.15), y + Inches(0.6), Inches(3.6), Inches(0.5),
        size=17, bold=True, color=col)
    for j, point in enumerate(points):
        txt(sl, f"• {point}", x + Inches(0.2), y + Inches(1.2 + j * 1.1),
            Inches(3.5), Inches(1.0), size=12, color=RGBColor(0xBB, 0xCC, 0xEE))

txt(sl, "Vos dossiers existants deviennent le moteur de l'IA — aucune saisie from scratch",
    Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.45),
    size=14, bold=True, italic=True, color=VERT, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Modèle commercial
# ═════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fond_bleu(sl)
barre_titre(sl, "Modèle commercial", "Licence logicielle par site / catégorie")

offres = [
    ("Starter", "1 catégorie\nJusqu'à 3 utilisateurs\nBase jusqu'à 50 références\nSupport email", BLEU_VIF, False),
    ("Pro", "Jusqu'à 5 catégories\nUtilisateurs illimités\nBase illimitée\nSupport prioritaire\nFormation incluse", VERT, True),
    ("Entreprise", "Catégories illimitées\nDéploiement sur site\nIntégration ERP possible\nSLA garanti\nCSM dédié", ORANGE, False),
]

for i, (nom, features, col, star) in enumerate(offres):
    x = Inches(0.5 + i * 4.2)
    y = Inches(1.6)
    h = Inches(4.8)
    rect(sl, x, y, Inches(3.9), Inches(0.55), col)
    rect(sl, x, y + Inches(0.55), Inches(3.9), h - Inches(0.55), RGBColor(0x12, 0x22, 0x44))
    txt(sl, nom, x, y, Inches(3.9), Inches(0.55),
        size=18, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    if star:
        txt(sl, "★ RECOMMANDÉ", x, y - Inches(0.4), Inches(3.9), Inches(0.35),
            size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
    for j, line in enumerate(features.split("\n")):
        txt(sl, f"✓  {line}", x + Inches(0.2), y + Inches(0.7 + j * 0.67),
            Inches(3.5), Inches(0.6), size=13, color=RGBColor(0xCC, 0xDD, 0xFF))

txt(sl, "Tarification sur devis selon volume de références et nombre de catégories",
    Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.45),
    size=13, italic=True, color=RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Contact / CTA
# ═════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fond_bleu(sl)
rect(sl, 0, 0, Inches(0.2), H, BLEU_VIF)
rect(sl, W - Inches(0.2), 0, Inches(0.2), H, BLEU_VIF)

txt(sl, "Prêt à diviser par 10\nle temps de vos dossiers qualité ?",
    Inches(1), Inches(1.2), Inches(11.3), Inches(2.0),
    size=36, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

txt(sl, "Demandez une démonstration sur vos propres données",
    Inches(1), Inches(3.4), Inches(11.3), Inches(0.6),
    size=18, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

rect(sl, Inches(4.2), Inches(4.2), Inches(4.9), Inches(0.7), BLEU_VIF)
txt(sl, "Démo gratuite — sans engagement",
    Inches(4.2), Inches(4.2), Inches(4.9), Inches(0.7),
    size=16, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

txt(sl, "Qual.IA Industrielle",
    Inches(1), Inches(5.8), Inches(11.3), Inches(0.5),
    size=14, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
txt(sl, "Système de génération automatique de dossiers qualité par intelligence artificielle",
    Inches(1), Inches(6.3), Inches(11.3), Inches(0.4),
    size=12, italic=True, color=RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.CENTER)

# ─── Sauvegarde ──────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"Presentation sauvegardee : {OUTPUT}")
